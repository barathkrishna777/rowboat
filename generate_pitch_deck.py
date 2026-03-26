"""Generate a pitch deck for Rowboat — CMU Startup Program application."""
# Branding: all references use "Rowboat" as the product name.

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ORANGE = RGBColor(0xFF, 0x69, 0x0F)
DARK = RGBColor(0x19, 0x20, 0x24)
GRAY = RGBColor(0x6B, 0x77, 0x85)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1D, 0xB9, 0x54)
LIGHT_BG = RGBColor(0xFA, 0xFB, 0xFC)
LIGHT_ORANGE_BG = RGBColor(0xFF, 0xF5, 0xED)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, corner_radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if corner_radius:
        shape.adjustments[0] = corner_radius
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_content(slide, left, top, width, bullets, font_size=18,
                              color=DARK, spacing=Pt(8), font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
    return txBox


def add_section_label(slide, text):
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(3), Inches(0.4),
                 text, font_size=12, color=ORANGE, bold=True)


def add_slide_title(slide, text):
    add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
                 text, font_size=36, color=DARK, bold=True)


def add_subtitle(slide, text, top=Inches(1.5)):
    add_text_box(slide, Inches(0.8), top, Inches(10), Inches(0.6),
                 text, font_size=18, color=GRAY)


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ── SLIDE 1: TITLE ────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, WHITE)

    add_shape_bg(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(4.2), DARK)

    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(1.0),
                 "ROWBOAT", font_size=54, color=ORANGE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(10), Inches(0.8),
                 "AI-Powered Group Outing Coordination", font_size=30, color=WHITE, bold=False)
    add_text_box(slide, Inches(0.8), Inches(2.8), Inches(10), Inches(0.6),
                 "From \"let's do something\" to a booked itinerary — in under a minute.",
                 font_size=18, color=RGBColor(0xA0, 0xA8, 0xB0))

    add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.5),
                 "CMU Startup Program Application", font_size=16, color=GRAY)
    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
                 "Carnegie Mellon University  •  24-880 AI Agents for Engineers",
                 font_size=14, color=GRAY)

    # ── SLIDE 2: PROBLEM ──────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_label(slide, "THE PROBLEM")
    add_slide_title(slide, "Planning a group outing is painful")

    problems = [
        "📱  47-message group chats just to pick a restaurant",
        "📅  Endless back-and-forth to find a time that works for everyone",
        "🤷  Decision fatigue — too many options, no one wants to decide",
        "🔀  Juggling 5+ apps: group chat, calendars, maps, Yelp, booking sites",
        "❌  Plans fall apart because coordination takes too long",
    ]
    add_bullet_slide_content(slide, Inches(0.8), Inches(2.2), Inches(11),
                              problems, font_size=22, spacing=Pt(14))

    add_shape_bg(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0), LIGHT_ORANGE_BG, 0.05)
    add_text_box(slide, Inches(1.2), Inches(5.95), Inches(11), Inches(0.6),
                 "Result: 60% of planned outings never happen because the coordination burden is too high.",
                 font_size=18, color=ORANGE, bold=True)

    # ── SLIDE 3: SOLUTION ─────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_label(slide, "THE SOLUTION")
    add_slide_title(slide, "One platform. One conversation. Done.")

    add_subtitle(slide, "Rowboat coordinates everything so your group doesn't have to.")

    features = [
        ("🎯", "Describe your outing in plain English", "\"Bowling night with dinner for 6 people next Saturday\""),
        ("🤖", "AI agents handle the rest", "Search venues, check calendars, match preferences, rank options"),
        ("👥", "Built for groups", "Friends list, one-click invites, shared preferences and history"),
        ("📅", "Book in one click", "Calendar invites sent to everyone automatically"),
    ]

    for i, (icon, title, desc) in enumerate(features):
        col_x = Inches(0.8 + (i * 3.1))
        add_shape_bg(slide, col_x, Inches(2.6), Inches(2.8), Inches(3.6), LIGHT_BG, 0.03)
        add_text_box(slide, col_x + Inches(0.3), Inches(2.8), Inches(2.2), Inches(0.6),
                     icon, font_size=36, color=DARK, alignment=PP_ALIGN.LEFT)
        add_text_box(slide, col_x + Inches(0.3), Inches(3.5), Inches(2.2), Inches(0.7),
                     title, font_size=17, color=DARK, bold=True)
        add_text_box(slide, col_x + Inches(0.3), Inches(4.2), Inches(2.2), Inches(1.5),
                     desc, font_size=14, color=GRAY)

    # ── SLIDE 4: HOW IT WORKS ─────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_label(slide, "HOW IT WORKS")
    add_slide_title(slide, "Multi-Agent AI Architecture")

    agents = [
        ("Search Agent", "Queries Google Places, Yelp,\nEventbrite, Ticketmaster", ORANGE),
        ("Preference Agent", "Adaptive questionnaire that\nbuilds rich user profiles", RGBColor(0x4A, 0x90, 0xD9)),
        ("Calendar Agent", "Finds time slots where\neveryone is available", GREEN),
        ("Recommendation Agent", "Constraint solver + RAG\nfor ranked suggestions", RGBColor(0x9B, 0x59, 0xB6)),
        ("Orchestrator", "Coordinates all agents into\na single planning flow", DARK),
    ]

    for i, (name, desc, color) in enumerate(agents):
        col_x = Inches(0.5 + (i * 2.5))
        shape = add_shape_bg(slide, col_x, Inches(2.4), Inches(2.2), Inches(2.8), WHITE, 0.04)
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

        add_text_box(slide, col_x + Inches(0.2), Inches(2.6), Inches(1.8), Inches(0.5),
                     name, font_size=15, color=color, bold=True)
        add_text_box(slide, col_x + Inches(0.2), Inches(3.2), Inches(1.8), Inches(1.5),
                     desc, font_size=13, color=GRAY)

    add_text_box(slide, Inches(0.8), Inches(5.6), Inches(11), Inches(0.8),
                 "Powered by PydanticAI  •  Gemini 2.5 Flash  •  ChromaDB RAG  •  FastAPI + Streamlit",
                 font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 5: MARKET ───────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_label(slide, "TARGET MARKET")
    add_slide_title(slide, "Starting with campus, scaling to every group")

    segments = [
        ("🎓", "Beachhead: College Students",
         "Frequent group outings, digitally native,\nstrong network effects on campus.\n20M+ US college students."),
        ("💼", "Expansion: Young Professionals",
         "Team dinners, friend groups in cities,\nrecurring social planning needs.\n44M Americans aged 22–34."),
        ("🏢", "Long-term: Organizations",
         "Student orgs, corporate teams,\nevent planning as a service.\nB2B SaaS opportunity."),
    ]

    for i, (icon, title, desc) in enumerate(segments):
        col_x = Inches(0.8 + (i * 4.0))
        add_shape_bg(slide, col_x, Inches(2.2), Inches(3.6), Inches(3.8), LIGHT_BG, 0.04)
        add_text_box(slide, col_x + Inches(0.3), Inches(2.4), Inches(3.0), Inches(0.6),
                     icon, font_size=40, color=DARK)
        add_text_box(slide, col_x + Inches(0.3), Inches(3.1), Inches(3.0), Inches(0.6),
                     title, font_size=18, color=DARK, bold=True)
        add_text_box(slide, col_x + Inches(0.3), Inches(3.8), Inches(3.0), Inches(2.0),
                     desc, font_size=15, color=GRAY)

    # ── SLIDE 6: BUSINESS MODEL ───────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_label(slide, "BUSINESS MODEL")
    add_slide_title(slide, "How we make money")

    models = [
        ("B2C: Freemium Subscriptions",
         "Free tier for basic planning.\nPremium unlocks smart personalization,\ncalendar sync, group memory,\nand priority booking."),
        ("B2B: Venue Partnerships",
         "Restaurants, entertainment venues,\nand franchises pay for featured\nplacement, lead generation,\nand booking conversions."),
        ("B2B: Organizational Licensing",
         "Universities, employers, and\nstudent orgs license Rowboat\nfor event coordination\nand social programming."),
    ]

    for i, (title, desc) in enumerate(models):
        col_x = Inches(0.8 + (i * 4.0))
        shape = add_shape_bg(slide, col_x, Inches(2.2), Inches(3.6), Inches(3.5), WHITE, 0.04)
        shape.line.color.rgb = RGBColor(0xE4, 0xE8, 0xEC)
        shape.line.width = Pt(1)

        add_text_box(slide, col_x + Inches(0.3), Inches(2.5), Inches(3.0), Inches(0.6),
                     title, font_size=17, color=ORANGE, bold=True)
        add_text_box(slide, col_x + Inches(0.3), Inches(3.2), Inches(3.0), Inches(2.2),
                     desc, font_size=15, color=GRAY)

    add_shape_bg(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.8), LIGHT_ORANGE_BG, 0.04)
    add_text_box(slide, Inches(1.2), Inches(6.1), Inches(11), Inches(0.6),
                 "Distribution: Product-led growth via group invites — every user brings in their friends.",
                 font_size=16, color=ORANGE, bold=True)

    # ── SLIDE 7: TRACTION ─────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_label(slide, "TRACTION & PROGRESS")
    add_slide_title(slide, "What we've built so far")

    milestones = [
        ("✅", "Working MVP", "End-to-end prototype with all 5 AI\nagents, constraint solver, and RAG pipeline"),
        ("✅", "Deployed", "Live on Railway with FastAPI backend\nand Streamlit frontend"),
        ("✅", "69 Tests", "Comprehensive test suite covering\nconstraints, models, DB, agents, and tools"),
        ("✅", "Social Features", "Friends system with requests,\nquick-add to groups, and user profiles"),
    ]

    for i, (icon, title, desc) in enumerate(milestones):
        col_x = Inches(0.5 + (i * 3.1))
        add_shape_bg(slide, col_x, Inches(2.2), Inches(2.9), Inches(3.0), LIGHT_BG, 0.04)
        add_text_box(slide, col_x + Inches(0.3), Inches(2.4), Inches(2.4), Inches(0.5),
                     f"{icon}  {title}", font_size=17, color=GREEN, bold=True)
        add_text_box(slide, col_x + Inches(0.3), Inches(3.1), Inches(2.4), Inches(1.8),
                     desc, font_size=14, color=GRAY)

    add_text_box(slide, Inches(0.8), Inches(5.6), Inches(11), Inches(0.8),
                 "Tech Stack: PydanticAI  •  Gemini 2.5 Flash  •  FastAPI  •  Streamlit  •  SQLAlchemy  •  ChromaDB",
                 font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 8: COMPETITIVE LANDSCAPE ────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_label(slide, "COMPETITIVE LANDSCAPE")
    add_slide_title(slide, "Why not just use what exists?")

    competitors = [
        ("Group Chats\n(iMessage, WhatsApp)", "No structure, no memory,\nno integration with venues\nor calendars.", "❌"),
        ("Yelp / Google Maps", "Great for discovery, but\nno group coordination,\npreferences, or scheduling.", "❌"),
        ("When2Meet / Doodle", "Scheduling only. No venue\ndiscovery, preferences,\nor booking.", "❌"),
        ("Rowboat", "All-in-one: preferences,\nschedules, discovery, ranking,\nbooking — AI-coordinated.", "✅"),
    ]

    for i, (name, desc, status) in enumerate(competitors):
        col_x = Inches(0.5 + (i * 3.1))
        bg_color = LIGHT_ORANGE_BG if status == "✅" else LIGHT_BG
        border_color = ORANGE if status == "✅" else RGBColor(0xE4, 0xE8, 0xEC)
        shape = add_shape_bg(slide, col_x, Inches(2.2), Inches(2.9), Inches(3.5), bg_color, 0.04)
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2 if status == "✅" else 1)

        add_text_box(slide, col_x + Inches(0.3), Inches(2.5), Inches(2.4), Inches(0.9),
                     name, font_size=15, color=DARK, bold=True)
        add_text_box(slide, col_x + Inches(0.3), Inches(3.5), Inches(2.4), Inches(1.5),
                     desc, font_size=14, color=GRAY)
        status_color = GREEN if status == "✅" else RGBColor(0xDC, 0x26, 0x26)
        label = "Full solution" if status == "✅" else "Partial"
        add_text_box(slide, col_x + Inches(0.3), Inches(5.0), Inches(2.4), Inches(0.4),
                     f"{status}  {label}", font_size=14, color=status_color, bold=True)

    # ── SLIDE 9: THE ASK ──────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_label(slide, "THE ASK")
    add_slide_title(slide, "What we need from the CMU Startup Program")

    asks = [
        "🧭  Mentorship on go-to-market strategy and customer discovery",
        "🏫  Access to the CMU community for early user validation on campus",
        "🤝  Connections to Pittsburgh's restaurant and entertainment ecosystem",
        "💡  Guidance on business model refinement and fundraising readiness",
        "🛠️  Resources to move from prototype to production-grade product",
    ]
    add_bullet_slide_content(slide, Inches(0.8), Inches(2.2), Inches(11),
                              asks, font_size=22, spacing=Pt(16))

    # ── SLIDE 10: CLOSING ─────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)

    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.0),
                 "ROWBOAT", font_size=54, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(2.8), Inches(11), Inches(0.8),
                 "Stop planning. Start going.",
                 font_size=28, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
                 "Carnegie Mellon University  •  24-880 AI Agents for Engineers",
                 font_size=16, color=RGBColor(0x8E, 0x99, 0xA4), alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(0.5),
                 "Thank you.",
                 font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # ── Save ──────────────────────────────────────────────────────────
    output_path = "Rowboat_Pitch_Deck.pptx"
    prs.save(output_path)
    print(f"Pitch deck saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    build_deck()
